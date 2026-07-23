# -*- coding: utf-8 -*-
"""
生成会议PPT: 抗干扰语音指令识别系统
要求: 字体大且清晰, 明亮风格
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== 配色方案 (明亮风格) =====
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE  = RGBColor(0x1F, 0x4E, 0x79)   # 主标题深蓝
BRIGHT_BLUE= RGBColor(0x2E, 0x9B, 0xD6)   # 亮蓝强调
ORANGE     = RGBColor(0xED, 0x7D, 0x31)   # 亮橙强调
GREEN      = RGBColor(0x2C, 0xA8, 0x6E)   # 亮绿
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)   # 浅灰区块
MID_GRAY   = RGBColor(0xBF, 0xBF, 0xBF)
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)   # 正文深灰
LIGHT_BLUE = RGBColor(0xDA, 0xE8, 0xF5)   # 浅蓝背景
LIGHT_ORG  = RGBColor(0xFD, 0xE9, 0xD9)   # 浅橙背景
LIGHT_GRN  = RGBColor(0xD5, 0xEA, 0xDB)   # 浅绿背景

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]

# ===== 辅助函数 =====
def add_bg(slide, color=WHITE):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, size=22, color=DARK_TEXT,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="微软雅黑"):
    """添加文本框"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return tb

def add_bullets(slide, left, top, width, height, items, size=22, color=DARK_TEXT,
                bullet_color=BRIGHT_BLUE, line_spacing=1.3):
    """添加项目符号列表"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # bullet符号
        r1 = p.add_run()
        r1.text = "● "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r1.font.name = "微软雅黑"
        # 内容
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "微软雅黑"
    return tb

def add_title_bar(slide, title, subtitle=None):
    """页面顶部标题栏 (明亮风格)"""
    # 顶部装饰条
    add_rect(slide, 0, 0, SW, Inches(1.15), DARK_BLUE)
    # 左侧亮色强调块
    add_rect(slide, 0, 0, Inches(0.18), Inches(1.15), ORANGE)
    # 标题文字
    add_text(slide, Inches(0.45), Inches(0.12), SW - Inches(2.5), Inches(0.7),
             title, size=32, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.72), SW - Inches(2.5), Inches(0.4),
                 subtitle, size=16, color=RGBColor(0xCC, 0xDD, 0xEE))
    # 右下角页码占位
    return

def add_footer(slide, page_num):
    """页脚"""
    add_rect(slide, 0, SH - Inches(0.32), SW, Inches(0.32), LIGHT_GRAY)
    add_text(slide, Inches(0.3), SH - Inches(0.32), Inches(8), Inches(0.32),
             "抗干扰语音指令识别系统 | XH-202615", size=11,
             color=RGBColor(0x88,0x88,0x88), anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(1.2), SH - Inches(0.32), Inches(0.9), Inches(0.32),
             str(page_num), size=11, color=RGBColor(0x88,0x88,0x88),
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def info_card(slide, left, top, width, height, title, lines, header_color=BRIGHT_BLUE,
              bg_color=LIGHT_BLUE, title_size=22, body_size=18):
    """信息卡片: 标题条 + 内容区"""
    # 内容背景
    add_rect(slide, left, top, width, height, bg_color)
    # 标题条
    add_rect(slide, left, top, width, Inches(0.55), header_color)
    add_text(slide, left + Inches(0.15), top, width - Inches(0.3), Inches(0.55),
             title, size=title_size, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # 内容
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.62),
                                   width - Inches(0.3), height - Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(body_size)
        r.font.color.rgb = DARK_TEXT
        r.font.name = "微软雅黑"
    return

# ============================================================
# 第1页: 封面
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
# 顶部装饰
add_rect(s, 0, 0, SW, Inches(0.35), DARK_BLUE)
add_rect(s, 0, Inches(0.35), SW, Inches(0.08), ORANGE)
# 主标题
add_text(s, Inches(0.8), Inches(1.8), SW - Inches(1.6), Inches(1.2),
         "复杂交互场景的\n抗干扰语音指令识别系统", size=48, color=DARK_BLUE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 副标题
add_text(s, Inches(0.8), Inches(3.4), SW - Inches(1.6), Inches(0.7),
         "V4.1 流水线技术报告", size=30, color=ORANGE, bold=True,
         align=PP_ALIGN.CENTER)
# 分隔线
add_rect(s, Inches(4.5), Inches(4.3), Inches(4.33), Inches(0.04), BRIGHT_BLUE)
# 信息
add_text(s, Inches(0.8), Inches(4.6), SW - Inches(1.6), Inches(0.5),
         "挑战杯 XH-202615 | 发榜单位: 美的集团", size=22, color=DARK_TEXT,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.2), SW - Inches(1.6), Inches(0.5),
         "降噪 → 人声分离 → 声纹鉴别 → 语音识别", size=24, color=BRIGHT_BLUE, bold=True,
         align=PP_ALIGN.CENTER)
# 底部
add_rect(s, 0, SH - Inches(0.5), SW, Inches(0.5), DARK_BLUE)
add_text(s, Inches(0.5), SH - Inches(0.5), SW - Inches(1), Inches(0.5),
         "团队会议汇报材料", size=16, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 第2页: 目录
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "汇报目录", "CONTENTS")
items = [
    ("01", "系统架构总览", "四阶段流水线与数据流"),
    ("02", "各模型输入输出详解", "降噪 / 分离 / 声纹 / ASR"),
    ("03", "系统使用方法", "环境配置与推理运行"),
    ("04", "比赛重要信息", "评分规则 / 数据集 / 时间节点"),
    ("05", "V4.1 性能成果", "当前最优结果对比"),
]
y = Inches(1.6)
for num, title, desc in items:
    # 序号圆块
    add_rect(s, Inches(1.0), y, Inches(0.8), Inches(0.8), ORANGE)
    add_text(s, Inches(1.0), y, Inches(0.8), Inches(0.8), num, size=26,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 标题与描述
    add_text(s, Inches(2.1), y + Inches(0.02), Inches(8), Inches(0.45),
             title, size=24, color=DARK_BLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(2.1), y + Inches(0.45), Inches(8), Inches(0.35),
             desc, size=16, color=RGBColor(0x66,0x66,0x66))
    y += Inches(1.05)
add_footer(s, 2)

# ============================================================
# 第3页: 系统架构总览
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "01  系统架构总览", "四阶段流水线 + 声纹决策中枢")

# 数据流说明
add_text(s, Inches(0.5), Inches(1.4), SW - Inches(1), Inches(0.5),
         "主阶段顺序: 降噪 → 人声分离 → 声纹鉴别 → 语音识别", size=22,
         color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

# 四个阶段方块
stages = [
    ("① 降噪", "noisereduce", "去除静态噪声", LIGHT_BLUE, BRIGHT_BLUE),
    ("② 人声分离", "SepFormer-16k", "盲分离多说话人", LIGHT_ORG, ORANGE),
    ("③ 声纹鉴别", "CAM++", "目标说话人验证", LIGHT_GRN, GREEN),
    ("④ 语音识别", "Paraformer", "文本转写", LIGHT_BLUE, BRIGHT_BLUE),
]
x = Inches(0.55)
w = Inches(2.85)
gap = Inches(0.35)
for name, model, desc, bg, hdr in stages:
    add_rect(s, x, Inches(2.15), w, Inches(1.9), bg)
    add_rect(s, x, Inches(2.15), w, Inches(0.55), hdr)
    add_text(s, x, Inches(2.15), w, Inches(0.55), name, size=20, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, Inches(2.8), w, Inches(0.45), model, size=16, color=DARK_BLUE,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.3), w, Inches(0.7), desc, size=15, color=DARK_TEXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    x += w + gap

# 箭头提示
add_text(s, Inches(0.5), Inches(4.15), SW - Inches(1), Inches(0.4),
         "→  →  →  →", size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# 核心设计要点
add_rect(s, Inches(0.5), Inches(4.7), SW - Inches(1), Inches(2.3), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.8), SW - Inches(1.4), Inches(0.5),
         "核心设计: 声纹是贯穿全流程的「决策中枢」", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(5.35), SW - Inches(1.6), Inches(1.6), [
    "声纹介入3次: ①kws提参考声纹  ②降噪音频算相似度决策是否分离  ③分离后选轨",
    "自适应分离: 仅当声纹相似度 < 0.28 (会被拒识) 时才触发分离, 避免污染已正确样本",
    "双阈值鉴别: 未分离样本阈值 0.28, 分离过样本阈值 0.35 (拦截neg假接受)",
], size=17, bullet_color=ORANGE, line_spacing=1.2)
add_footer(s, 3)

# ============================================================
# 第4页: 贯穿数据契约
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "数据契约: 所有模型共用一条「总线」", "THE UNIFIED INTERFACE")

# 核心契约
add_rect(s, Inches(0.5), Inches(1.5), SW - Inches(1), Inches(1.4), LIGHT_BLUE)
add_text(s, Inches(0.7), Inches(1.6), SW - Inches(1.4), Inches(0.5),
         "音频格式 (所有模型的统一输入输出)", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(2.1), SW - Inches(1.6), Inches(0.8), [
    "np.ndarray  |  float32  |  取值范围 [-1, 1]  |  单声道  |  采样率 16 kHz",
], size=20, bullet_color=BRIGHT_BLUE, line_spacing=1.1)

# 声纹格式
add_rect(s, Inches(0.5), Inches(3.05), SW - Inches(1), Inches(1.0), LIGHT_GRN)
add_text(s, Inches(0.7), Inches(3.12), SW - Inches(1.4), Inches(0.45),
         "声纹格式", size=20, color=GREEN, bold=True)
add_bullets(s, Inches(0.8), Inches(3.55), SW - Inches(1.6), Inches(0.5), [
    "L2 归一化的一维向量 (CAM++ 为 192 维)  |  用 cosine 相似度比较",
], size=19, bullet_color=GREEN, line_spacing=1.1)

# 模块对接方式
add_rect(s, Inches(0.5), Inches(4.2), SW - Inches(1), Inches(2.7), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.3), SW - Inches(1.4), Inches(0.5),
         "模块对接方式", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(4.85), SW - Inches(1.6), Inches(2.0), [
    "每个模块继承 Base 基类, 实现 load() + 一个核心方法",
    "核心方法: denoise / separate / extract / transcribe",
    "工厂函数 create_xxx(config, device) 根据 config.model 字段实例化具体实现",
    "换模型 = 写一个新子类 + 在 config 里改 model 字段 (pipeline 逻辑不动)",
], size=18, bullet_color=ORANGE, line_spacing=1.25)
add_footer(s, 4)

# ============================================================
# 第5页: 模型①降噪 noisereduce
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "模型①  noisereduce 降噪", "DENOISER MODULE")

# 接口
info_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.2),
          "接口定义",
          ["denoise(audio, sr=16000) → np.ndarray",
           "",
           "输入: cmd 原始音频 (float32, [-1,1], 16kHz)",
           "输出: 降噪后音频 (同格式, 去除非人声噪声)",
           "执行时机: 每个样本必跑, 无条件"],
          header_color=BRIGHT_BLUE, bg_color=LIGHT_BLUE, title_size=22, body_size=17)

# 模型信息
info_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.2),
          "当前模型",
          ["算法: 频谱门限 (非深度学习)",
           "配置: stationary=True, prop_decrease=0.8",
           "特点: 极轻量, 去静态噪声",
           "局限: 无法区分目标人声与干扰人声"],
          header_color=ORANGE, bg_color=LIGHT_ORG, title_size=22, body_size=17)

# 对接要点
add_rect(s, Inches(0.5), Inches(3.95), SW - Inches(1), Inches(3.0), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.05), SW - Inches(1.4), Inches(0.5),
         "对接要点 & 更换注意事项", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(4.6), SW - Inches(1.6), Inches(2.3), [
    "降噪是流水线「必经之路」, 输出直接喂给声纹提取和分离",
    "效率敏感: 每样本必跑, 别换太重的模型",
    "换 DeepFilterNet3: 原生48kHz需16k↔48k重采样, 可能引入质量损失; 需Rust编译",
    "换 FullSubNet+: 需手动下载checkpoint, STFT分帧处理",
    "比赛约束: 模型权重需提前下载到 pretrained/, 运行时不联网",
], size=17, bullet_color=ORANGE, line_spacing=1.2)
add_footer(s, 5)

# ============================================================
# 第6页: 模型②SepFormer人声分离
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "模型②  SepFormer-16k 人声分离", "SEPARATOR MODULE")

info_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.4),
          "接口定义",
          ["separate(audio, sr=16000, target_embedding=None)",
           "  → (best_audio, all_sources[])",
           "",
           "输入: 降噪后音频 (float32, [-1,1], 16kHz)",
           "输出: 元组 (最佳音轨, 全部音轨列表)",
           "触发: 仅当 sim_denoised < 0.28 时调用"],
          header_color=ORANGE, bg_color=LIGHT_ORG, title_size=21, body_size=16)

info_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.4),
          "当前模型",
          ["模型: speechbrain/sepformer-whamr16k",
           "类型: 16kHz原生盲分离",
           "原理: 无差别分开所有声源",
           "输出: 2条音轨 (比赛最多2说话人)",
           "选轨: pipeline层用CAM++声纹选轨"],
          header_color=BRIGHT_BLUE, bg_color=LIGHT_BLUE, title_size=21, body_size=16)

# 重要坑
add_rect(s, Inches(0.5), Inches(4.1), SW - Inches(1), Inches(2.85), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.2), SW - Inches(1.4), Inches(0.5),
         "重要坑 & 更换注意事项", size=22, color=RGBColor(0xC0,0x39,0x2B), bold=True)
add_bullets(s, Inches(0.8), Inches(4.75), SW - Inches(1.6), Inches(2.1), [
    "separator.py 内部 _select_best_match 用能量法选轨 (选声音最大的) → V3踩坑根源",
    "V4修复: 在 pipeline.py 层做声纹选轨, 对每条音轨提声纹选与kws最相似的",
    "绝不用 SepFormer-8k: 8kHz降采样导致高频丢失, V2已验证会降质",
    "换 SpEx+ (目标说话人提取): 输出直接是目标人声不需选轨, pipeline逻辑要改",
    "分离最耗时 (V4.1占1064s绝大部分), 换模型重点看推理速度",
], size=16, bullet_color=RGBColor(0xC0,0x39,0x2B), line_spacing=1.15)
add_footer(s, 6)

# ============================================================
# 第7页: 模型③CAM++声纹鉴别
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "模型③  CAM++ 声纹鉴别", "VOICEPRINT MODULE")

info_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.4),
          "接口定义",
          ["extract(audio, sr=16000) → np.ndarray [192]",
           "",
           "输入: 音频 (float32, [-1,1], 16kHz)",
           "输出: 192维 L2归一化 embedding 向量",
           "决策: cosine相似度 + 阈值判断 (pipeline层)"],
          header_color=GREEN, bg_color=LIGHT_GRN, title_size=21, body_size=16)

info_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.4),
          "当前模型 & 调用",
          ["模型: iic/speech_campplus_sv_zh-cn_16k-common",
           "来源: ModelScope",
           "",
           "每样本调用3次:",
           "① kws唤醒音频 → 参考声纹",
           "② 降噪音频 → 算 sim_denoised",
           "③ 分离后每条音轨 → 选轨"],
          header_color=BRIGHT_BLUE, bg_color=LIGHT_BLUE, title_size=21, body_size=16)

add_rect(s, Inches(0.5), Inches(4.1), SW - Inches(1), Inches(2.85), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.2), SW - Inches(1.4), Inches(0.5),
         "更换注意事项", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(4.75), SW - Inches(1.6), Inches(2.1), [
    "维度可能变: CAM++/ECAPA-TDNN为192维, WeSpeaker等可能不同 (cosine不依赖维度)",
    "阈值必须重新校准: 双阈值0.28/0.35是CAM++专属, 换模型需在datasetA重新扫描",
    "效率敏感: 每样本调用3次, 是被调用最频繁的模块",
    "兼容性: voiceprint.py顶部有PyTorch2.5 FSDP2兼容patch, 遇CPUOffloadPolicy报错确认还在",
    "声纹模块只产出embedding, 接受/拒识判断都在pipeline层 → 换模型pipeline逻辑不动",
], size=16, bullet_color=GREEN, line_spacing=1.15)
add_footer(s, 7)

# ============================================================
# 第8页: 模型④Paraformer ASR
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "模型④  Paraformer 语音识别", "ASR MODULE")

info_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.4),
          "接口定义",
          ["transcribe(audio, sr=16000) → str",
           "",
           "输入: best_audio (经normalize+trim预处理)",
           "输出: 识别文本字符串 (含标点)",
           "后处理: pipeline层 strip_punctuation 去标点",
           "执行: 仅对通过鉴别的样本调用"],
          header_color=BRIGHT_BLUE, bg_color=LIGHT_BLUE, title_size=21, body_size=16)

info_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.4),
          "当前模型",
          ["主模型: paraformer-zh (FunASR)",
           "VAD: fsmn-vad (自动断句)",
           "标点: ct-punc (标点恢复)",
           "热词: 配置领域热词 (空调/洗碗机等)",
           "输入形式: 需临时WAV文件 (非numpy)"],
          header_color=ORANGE, bg_color=LIGHT_ORG, title_size=21, body_size=16)

add_rect(s, Inches(0.5), Inches(4.1), SW - Inches(1), Inches(2.85), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.2), SW - Inches(1.4), Inches(0.5),
         "更换注意事项", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(4.75), SW - Inches(1.6), Inches(2.1), [
    "输入形式差异: Paraformer需临时WAV文件; Whisper/SherpaONNX可直接接受numpy",
    "标点处理: pipeline层统一strip_punctuation, 模型不加标点也无害 (空操作)",
    "热词能力: 换Whisper/SherpaONNX会丢失领域热词, 影响专业词汇识别",
    "ASR最贵: 只对通过鉴别的样本跑, 拿到最干净音轨时CER最低",
    "比赛标签不含标点 → ASR输出必须去标点后再算CER",
], size=17, bullet_color=BRIGHT_BLUE, line_spacing=1.2)
add_footer(s, 8)

# ============================================================
# 第9页: 系统使用方法 - 环境
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "03  系统使用方法 (1/2)", "环境配置与依赖安装")

# 步骤1
add_rect(s, Inches(0.5), Inches(1.5), Inches(0.7), Inches(0.7), ORANGE)
add_text(s, Inches(0.5), Inches(1.5), Inches(0.7), Inches(0.7), "1", size=28,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.4), Inches(1.5), Inches(11), Inches(0.7),
         "克隆仓库", size=22, color=DARK_BLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.4), Inches(2.15), Inches(11), Inches(0.4),
         "git clone https://github.com/Wowater404/Anti-interference-voice-system.git",
         size=16, color=DARK_TEXT, font="Consolas")

# 步骤2
add_rect(s, Inches(0.5), Inches(2.75), Inches(0.7), Inches(0.7), ORANGE)
add_text(s, Inches(0.5), Inches(2.75), Inches(0.7), Inches(0.7), "2", size=28,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.4), Inches(2.75), Inches(11), Inches(0.7),
         "安装依赖", size=22, color=DARK_BLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.4), Inches(3.4), Inches(11), Inches(0.4),
         "pip install -r requirements.txt", size=16, color=DARK_TEXT, font="Consolas")

# 步骤3
add_rect(s, Inches(0.5), Inches(4.0), Inches(0.7), Inches(0.7), ORANGE)
add_text(s, Inches(0.5), Inches(4.0), Inches(0.7), Inches(0.7), "3", size=28,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.4), Inches(4.0), Inches(11), Inches(0.7),
         "首次运行自动下载模型权重", size=22, color=DARK_BLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(1.4), Inches(4.65), Inches(11), Inches(0.4),
         "CAM++ (ModelScope) / SepFormer (HuggingFace) / Paraformer (FunASR)",
         size=15, color=DARK_TEXT, font="Consolas")

# 关键依赖
add_rect(s, Inches(0.5), Inches(5.3), SW - Inches(1), Inches(1.65), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(5.38), SW - Inches(1.4), Inches(0.45),
         "关键依赖", size=20, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(5.8), SW - Inches(1.6), Inches(1.1), [
    "PyTorch 2.5+ (CUDA版用于GPU加速)  |  FunASR  |  speechbrain  |  modelscope  |  noisereduce  |  librosa  |  pyyaml",
], size=16, bullet_color=ORANGE, line_spacing=1.1)
add_footer(s, 9)

# ============================================================
# 第10页: 系统使用方法 - 运行
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "03  系统使用方法 (2/2)", "推理运行与输出")

# 运行命令
add_rect(s, Inches(0.5), Inches(1.5), SW - Inches(1), Inches(1.8), LIGHT_BLUE)
add_text(s, Inches(0.7), Inches(1.58), SW - Inches(1.4), Inches(0.5),
         "运行推理", size=22, color=DARK_BLUE, bold=True)
add_text(s, Inches(0.8), Inches(2.1), SW - Inches(1.6), Inches(1.1),
         "python run_inference.py \\\n  --data_root <数据集路径> \\\n  --split all \\\n  --output results/output.json \\\n  --checkpoint results/checkpoint.json",
         size=17, color=DARK_TEXT, font="Consolas")

# 参数说明
info_card(s, Inches(0.5), Inches(3.5), Inches(6.0), Inches(3.4),
          "关键参数",
          ["--data_root: 数据集根目录 (含pos/和neg/)",
           "--split: all | pos | neg",
           "--output: 结果输出JSON路径",
           "--checkpoint: 断点续传文件路径",
           "",
           "配置文件: configs/default.yaml",
           "  separation.enable: 是否启用分离",
           "  voiceprint.threshold: 声纹阈值(0.28)",
           "  vp_threshold_separated: 分离后阈值(0.35)"],
          header_color=GREEN, bg_color=LIGHT_GRN, title_size=20, body_size=15)

# 输出格式
info_card(s, Inches(6.8), Inches(3.5), Inches(6.0), Inches(3.4),
          "输出格式 (比赛提交JSON)",
          ["{",
           "  \"result\": {",
           "    \"results\": [",
           "      {\"id\": \"...\", \"content\": \"识别文本\",",
           "       \"label\": \"标签\", \"cer\": \"0.xx\"}",
           "    ],",
           "    \"final_cer\": \"0.5738\",",
           "    \"duration\": \"1064.29\"",
           "  },",
           "  \"metrics\": {\"rejection_rate\": \"0.9367\", ...}",
           "}"],
          header_color=BRIGHT_BLUE, bg_color=LIGHT_BLUE, title_size=20, body_size=14)
add_footer(s, 10)

# ============================================================
# 第11页: 比赛重要信息(1) - 基本规则
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "04  比赛重要信息 (1/2)", "评分规则与时间节点")

# 评分公式
add_rect(s, Inches(0.5), Inches(1.5), SW - Inches(1), Inches(1.3), LIGHT_ORG)
add_text(s, Inches(0.7), Inches(1.58), SW - Inches(1.4), Inches(0.5),
         "评分公式", size=22, color=ORANGE, bold=True)
add_text(s, Inches(0.8), Inches(2.05), SW - Inches(1.6), Inches(0.7),
         "总分 = CER × 40%  +  RR × 40%  +  推理效率 × 20% (时间10% + 内存10%)",
         size=22, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

# 三大指标
info_card(s, Inches(0.5), Inches(3.0), Inches(4.0), Inches(1.8),
          "CER (字错误率)",
          ["pos样本的字错误率",
           "越低越好 (理想=0)",
           "拒识样本按删除错误算",
           "权重40%"],
          header_color=BRIGHT_BLUE, bg_color=LIGHT_BLUE, title_size=19, body_size=15)

info_card(s, Inches(4.65), Inches(3.0), Inches(4.0), Inches(1.8),
          "RR (拒识率)",
          ["neg样本被正确拒识的比例",
           "越高越好 (理想=1)",
           "neg只算RR不算CER",
           "权重40%"],
          header_color=ORANGE, bg_color=LIGHT_ORG, title_size=19, body_size=15)

info_card(s, Inches(8.8), Inches(3.0), Inches(4.0), Inches(1.8),
          "推理效率",
          ["推理时间 + 内存占用",
           "统一GPU环境: L20-46G",
           "需提交GPU推理脚本",
           "权重20%"],
          header_color=GREEN, bg_color=LIGHT_GRN, title_size=19, body_size=15)

# 时间节点
add_rect(s, Inches(0.5), Inches(5.0), SW - Inches(1), Inches(1.95), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(5.08), SW - Inches(1.4), Inches(0.5),
         "时间节点", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(5.6), SW - Inches(1.6), Inches(1.3), [
    "6/30: 报名截止   |   9/5: 提交模型 + GPU推理脚本 (datasetB最终排名依据)",
    "11月: 答辩   |   datasetA为临时排行榜, 仅开发参考, 不决定入围",
], size=18, bullet_color=ORANGE, line_spacing=1.3)
add_footer(s, 11)

# ============================================================
# 第12页: 比赛重要信息(2) - 数据集与CER
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "04  比赛重要信息 (2/2)", "数据集规格与CER计算")

# 数据集
info_card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.6),
          "数据集 datasetA",
          ["pos: 1364条 (测CER) | neg: 474条 (测RR)",
           "音频: 16kHz / mono / 16bit WAV",
           "kws唤醒词: 1.5s | cmd指令: 1.28~6.2s",
           "20种唤醒词(pos), neg多2种变体",
           "(hicolmo, 空调开机)",
           "datasetB: 最终排名依据 (9/5前提交)"],
          header_color=BRIGHT_BLUE, bg_color=LIGHT_BLUE, title_size=20, body_size=15)

# CER计算
info_card(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.6),
          "CER 计算 (官方标准)",
          ["1. NFKC Unicode归一化",
           "2. lowercase 转小写",
           "3. 全 Unicode P* 标点过滤",
           "4. editdistance 编辑距离",
           "5. micro-average: total_errors/total_chars",
           "拒识输出空字符串 \"\" (按删除错误算)"],
          header_color=ORANGE, bg_color=LIGHT_ORG, title_size=20, body_size=15)

# FAQ关键点
add_rect(s, Inches(0.5), Inches(4.3), SW - Inches(1), Inches(2.65), LIGHT_GRAY)
add_text(s, Inches(0.7), Inches(4.38), SW - Inches(1.4), Inches(0.5),
         "比赛关键 FAQ", size=22, color=DARK_BLUE, bold=True)
add_bullets(s, Inches(0.8), Inches(4.9), SW - Inches(1.6), Inches(2.0), [
    "datasetA是测评集不是训练集 → 模型为预训练直接推理 (zero-shot), 不能用它训练/微调",
    "9/5前提交: 模型权重 + GPU推理脚本, 不是提交代码",
    "统一环境L20-46G, 必须能用CUDA推理; 模型权重提前下载, 运行时不联网",
    "neg样本只算RR不算CER; pos错误拒识按删除错误计算CER (输出空字符串)",
], size=17, bullet_color=ORANGE, line_spacing=1.2)
add_footer(s, 12)

# ============================================================
# 第13页: V4.1 性能成果
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_title_bar(s, "05  V4.1 性能成果", "当前最优结果对比")

# 对比表头
add_rect(s, Inches(0.5), Inches(1.5), SW - Inches(1), Inches(0.6), DARK_BLUE)
headers = ["指标", "V2 (无分离)", "V4 (分离 t=0.28)", "V4.1 (双阈值) ✓"]
col_w = [Inches(2.5), Inches(3.1), Inches(3.1), Inches(3.13)]
x = Inches(0.5)
for h, w in zip(headers, col_w):
    add_text(s, x, Inches(1.5), w, Inches(0.6), h, size=18, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += w

# 数据行
rows = [
    ["CER (micro)", "0.5857", "0.5615", "0.5738"],
    ["RR", "0.9367", "0.9051", "0.9367"],
    ["Score (CER40+RR40)", "0.5404", "0.5374", "0.5452"],
    ["Pos 接受数", "900", "948", "913"],
    ["Neg 假接受数", "30", "45", "30"],
]
y = Inches(2.1)
for i, row in enumerate(rows):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    if i == 2:  # Score行高亮
        bg = LIGHT_GRN
    add_rect(s, Inches(0.5), y, SW - Inches(1), Inches(0.55), bg)
    x = Inches(0.5)
    for j, (val, w) in enumerate(zip(row, col_w)):
        color = DARK_TEXT
        bold = False
        if j == 0:
            color = DARK_BLUE; bold = True
        if j == 3:  # V4.1列
            color = GREEN; bold = True
        add_text(s, x, y, w, Inches(0.55), val, size=17, color=color, bold=bold,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += w
    y += Inches(0.55)

# 结论
add_rect(s, Inches(0.5), Inches(5.1), SW - Inches(1), Inches(1.85), LIGHT_GRN)
add_text(s, Inches(0.7), Inches(5.18), SW - Inches(1.4), Inches(0.5),
         "V4.1 结论", size=22, color=GREEN, bold=True)
add_bullets(s, Inches(0.8), Inches(5.7), SW - Inches(1.6), Inches(1.2), [
    "CER降低0.012, RR完全不变, Score提升+0.0048 → 成功让分离模块生效并超越V2基线",
    "救回13个pos样本, 0个neg假接受 (双阈值0.35完美拦截分离后的neg假接受)",
], size=17, bullet_color=GREEN, line_spacing=1.25)
add_footer(s, 13)

# ============================================================
# 第14页: 总结
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
add_rect(s, 0, 0, SW, Inches(0.35), DARK_BLUE)
add_rect(s, 0, Inches(0.35), SW, Inches(0.08), ORANGE)

add_text(s, Inches(0.8), Inches(1.2), SW - Inches(1.6), Inches(0.9),
         "总结与展望", size=40, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

# 三大要点
points = [
    ("完整流水线", "降噪→分离→声纹→ASR 四阶段, 声纹作为决策中枢贯穿全流程", BRIGHT_BLUE, LIGHT_BLUE),
    ("V4.1 双阈值", "自适应分离+声纹选轨+双阈值, CER 0.5738 / RR 0.9367 / Score 0.5452", ORANGE, LIGHT_ORG),
    ("模块可替换", "统一接口契约, 换模型只需写新子类+改config, pipeline逻辑不动", GREEN, LIGHT_GRN),
]
y = Inches(2.4)
for title, desc, hdr, bg in points:
    add_rect(s, Inches(1.0), y, SW - Inches(2), Inches(1.15), bg)
    add_rect(s, Inches(1.0), y, Inches(0.15), Inches(1.15), hdr)
    add_text(s, Inches(1.35), y + Inches(0.1), SW - Inches(2.5), Inches(0.45),
             title, size=22, color=hdr, bold=True)
    add_text(s, Inches(1.35), y + Inches(0.55), SW - Inches(2.5), Inches(0.55),
             desc, size=16, color=DARK_TEXT)
    y += Inches(1.3)

# 后续方向
add_rect(s, Inches(1.0), y, SW - Inches(2), Inches(1.1), LIGHT_GRAY)
add_text(s, Inches(1.2), y + Inches(0.05), SW - Inches(2.4), Inches(0.45),
         "后续优化方向", size=20, color=DARK_BLUE, bold=True)
add_text(s, Inches(1.2), y + Inches(0.5), SW - Inches(2.4), Inches(0.55),
         "SpEx+目标说话人提取  |  sep_trigger_min优化分离触发范围  |  L20-46G GPU推理加速",
         size=15, color=DARK_TEXT)

add_rect(s, 0, SH - Inches(0.5), SW, Inches(0.5), DARK_BLUE)
add_text(s, Inches(0.5), SH - Inches(0.5), SW - Inches(1), Inches(0.5),
         "谢谢观看  |  抗干扰语音指令识别系统  |  XH-202615", size=18, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===== 保存 =====
out = r"F:\龙虾\2026-07-18-13-57-00\voice_pipeline\抗干扰语音指令识别系统_会议PPT.pptx"
prs.save(out)
print(f"PPT已生成: {out}")
print(f"共 {len(prs.slides)} 页")
